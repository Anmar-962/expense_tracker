from pptx import Presentation
from PIL import Image
import pytesseract
import io

# Set the path to tesseract executable
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Load the PowerPoint file
prs = Presentation(r"C:\Users\MSI\Desktop\College\Fundamentals Of Databases\Chapter 01.pptx")  # Replace with your actual file

# Loop through slides and extract text from images
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # 13 means it's a picture/image
            image = shape.image
            image_bytes = image.blob
            img = Image.open(io.BytesIO(image_bytes))
            text = pytesseract.image_to_string(img)
            print(f"Slide {i+1} Text:\n{text}\n{'-'*40}")
